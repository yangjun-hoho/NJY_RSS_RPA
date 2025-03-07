import streamlit as st
import io
import tempfile
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import PyPDF2
import docx
import re
import time
import openai
from openai import OpenAI
import google.generativeai as genai
from dotenv import load_dotenv

# API 키 로드
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")

# 향상된 PPT 템플릿 정보
ENHANCED_TEMPLATES = {
    "기본 템플릿": {
        # 기본 색상 정보
        "colors": {
            "bg_color": (255, 255, 255),            # 배경색
            "title_color": (31, 73, 125),           # 제목 텍스트 색상
            "text_color": (0, 0, 0),                # 본문 텍스트 색상
            "accent_color": (79, 129, 189),         # 강조 색상
            "footer_color": (128, 128, 128),        # 푸터 텍스트 색상
            "bullet_color": (79, 129, 189),         # 불릿 색상
            "header_color": (31, 73, 125),          # 헤더 텍스트 색상
            "shape_fill": (230, 240, 250),          # 도형 채우기 색상
            "shape_line": (31, 73, 125),            # 도형 선 색상
        },
        # 폰트 정보
        "fonts": {
            "title_font": {
                "name": "맑은 고딕",
                "size": Pt(40),                     # 제목 폰트 크기
                "bold": True,
                "italic": False,
                "underline": False
            },
            "subtitle_font": {
                "name": "맑은 고딕",
                "size": Pt(24),                     # 부제목 폰트 크기
                "bold": False,
                "italic": False,
                "underline": False
            },
            "body_font": {
                "name": "맑은 고딕",
                "size": Pt(18),                     # 본문 폰트 크기
                "bold": False,
                "italic": False,
                "underline": False
            },
            "footer_font": {
                "name": "맑은 고딕",
                "size": Pt(10),                     # 푸터 폰트 크기
                "bold": False,
                "italic": False,
                "underline": False
            }
        },
        # 스타일 정보
        "styles": {
            "title_align": PP_ALIGN.CENTER,         # 제목 정렬
            "body_align": PP_ALIGN.LEFT,            # 본문 정렬
            "margins": {
                "top": Inches(1),
                "bottom": Inches(0.8),
                "left": Inches(0.8),
                "right": Inches(0.8)
            },
            "bullet_style": "•",                    # 불릿 스타일
            "slide_number": True,                  # 슬라이드 번호 표시 여부
            "shape_roundness": 0.1,                # 도형 모서리 둥글기 정도 (0-1)
        },
        # 기타 설정
        "settings": {
            "footer_text": "© 2025 남양주시",       # 기본 푸터 텍스트
            "include_logo": False,                 # 로고 포함 여부
            "logo_path": None,                     # 로고 경로
            "background_style": "solid",           # 배경 스타일 (solid, gradient, pattern)
            "gradient_stops": None,                # 그라데이션 지점 설정
        }
    },
    
    "모던 블루": {
        # 기본 색상 정보
        "colors": {
            "bg_color": (255, 255, 255),
            "title_color": (0, 112, 192),
            "text_color": (68, 68, 68),
            "accent_color": (0, 176, 240),
            "footer_color": (120, 120, 120),
            "bullet_color": (0, 176, 240),
            "header_color": (0, 112, 192),
            "shape_fill": (235, 245, 255),
            "shape_line": (0, 112, 192),
        },
        # 폰트 정보
        "fonts": {
            "title_font": {
                "name": "나눔스퀘어",
                "size": Pt(44),
                "bold": True,
                "italic": False,
                "underline": False
            },
            "subtitle_font": {
                "name": "나눔스퀘어",
                "size": Pt(28),
                "bold": False,
                "italic": False,
                "underline": False
            },
            "body_font": {
                "name": "나눔스퀘어",
                "size": Pt(20),
                "bold": False,
                "italic": False,
                "underline": False
            },
            "footer_font": {
                "name": "나눔스퀘어",
                "size": Pt(11),
                "bold": False,
                "italic": False,
                "underline": False
            }
        },
        # 스타일 정보
        "styles": {
            "title_align": PP_ALIGN.LEFT,
            "body_align": PP_ALIGN.LEFT,
            "margins": {
                "top": Inches(0.8),
                "bottom": Inches(0.6),
                "left": Inches(0.8),
                "right": Inches(0.6)
            },
            "bullet_style": "➤",
            "slide_number": True,
            "shape_roundness": 0.25,
        },
        # 기타 설정
        "settings": {
            "footer_text": "남양주시 | AI 생성 프레젠테이션",
            "include_logo": False,
            "logo_path": None,
            "background_style": "gradient",
            "gradient_stops": [(0, (255, 255, 255)), (1.0, (240, 250, 255))],
        }
    },
    
    "다크 테마": {
        # 기본 색상 정보
        "colors": {
            "bg_color": (47, 47, 47),
            "title_color": (255, 255, 255),
            "text_color": (235, 235, 235),
            "accent_color": (255, 89, 94),
            "footer_color": (180, 180, 180),
            "bullet_color": (255, 89, 94),
            "header_color": (255, 255, 255),
            "shape_fill": (70, 70, 70),
            "shape_line": (255, 89, 94),
        },
        # 폰트 정보
        "fonts": {
            "title_font": {
                "name": "맑은 고딕",
                "size": Pt(42),
                "bold": True,
                "italic": False,
                "underline": False
            },
            "subtitle_font": {
                "name": "맑은 고딕",
                "size": Pt(26),
                "bold": False,
                "italic": False,
                "underline": False
            },
            "body_font": {
                "name": "맑은 고딕",
                "size": Pt(20),
                "bold": False,
                "italic": False,
                "underline": False
            },
            "footer_font": {
                "name": "맑은 고딕",
                "size": Pt(10),
                "bold": False,
                "italic": False,
                "underline": False
            }
        },
        # 스타일 정보
        "styles": {
            "title_align": PP_ALIGN.LEFT,
            "body_align": PP_ALIGN.LEFT,
            "margins": {
                "top": Inches(0.8),
                "bottom": Inches(0.8),
                "left": Inches(1.0),
                "right": Inches(0.8)
            },
            "bullet_style": "–",
            "slide_number": True,
            "shape_roundness": 0.2,
        },
        # 기타 설정
        "settings": {
            "footer_text": "남양주시 © 2025",
            "include_logo": False,
            "logo_path": None,
            "background_style": "gradient",
            "gradient_stops": [(0, (47, 47, 47)), (1.0, (30, 30, 30))],
        }
    },
    
    "공식 보고서": {
        # 기본 색상 정보
        "colors": {
            "bg_color": (245, 245, 245),
            "title_color": (0, 51, 102),
            "text_color": (51, 51, 51),
            "accent_color": (192, 0, 0),
            "footer_color": (100, 100, 100),
            "bullet_color": (0, 51, 102),
            "header_color": (0, 51, 102),
            "shape_fill": (225, 230, 235),
            "shape_line": (0, 51, 102),
        },
        # 폰트 정보
        "fonts": {
            "title_font": {
                "name": "굴림",
                "size": Pt(36),
                "bold": True,
                "italic": False,
                "underline": False
            },
            "subtitle_font": {
                "name": "굴림",
                "size": Pt(22),
                "bold": True,
                "italic": False,
                "underline": False
            },
            "body_font": {
                "name": "굴림",
                "size": Pt(18),
                "bold": False,
                "italic": False,
                "underline": False
            },
            "footer_font": {
                "name": "굴림",
                "size": Pt(9),
                "bold": False,
                "italic": False,
                "underline": False
            }
        },
        # 스타일 정보
        "styles": {
            "title_align": PP_ALIGN.CENTER,
            "body_align": PP_ALIGN.LEFT,
            "margins": {
                "top": Inches(0.7),
                "bottom": Inches(0.7),
                "left": Inches(0.7),
                "right": Inches(0.7)
            },
            "bullet_style": "■",
            "slide_number": True,
            "shape_roundness": 0.05,
        },
        # 기타 설정
        "settings": {
            "footer_text": "남양주시 공식 보고서 | 보안등급: 일반",
            "include_logo": False,
            "logo_path": None,
            "background_style": "pattern",
            "pattern_style": "subtle_grid",
        }
    },
    
    "미니멀리즘": {
        # 기본 색상 정보
        "colors": {
            "bg_color": (255, 255, 255),
            "title_color": (0, 0, 0),
            "text_color": (80, 80, 80),
            "accent_color": (0, 150, 136),
            "footer_color": (150, 150, 150),
            "bullet_color": (0, 150, 136),
            "header_color": (0, 0, 0),
            "shape_fill": (240, 240, 240),
            "shape_line": (200, 200, 200),
        },
        # 폰트 정보
        "fonts": {
            "title_font": {
                "name": "맑은 고딕",
                "size": Pt(40),
                "bold": False,
                "italic": False,
                "underline": False
            },
            "subtitle_font": {
                "name": "맑은 고딕",
                "size": Pt(24),
                "bold": False,
                "italic": False,
                "underline": False
            },
            "body_font": {
                "name": "맑은 고딕",
                "size": Pt(18),
                "bold": False,
                "italic": False,
                "underline": False
            },
            "footer_font": {
                "name": "맑은 고딕",
                "size": Pt(10),
                "bold": False,
                "italic": False,
                "underline": False
            }
        },
        # 스타일 정보
        "styles": {
            "title_align": PP_ALIGN.LEFT,
            "body_align": PP_ALIGN.LEFT,
            "margins": {
                "top": Inches(1.0),
                "bottom": Inches(1.0),
                "left": Inches(1.2),
                "right": Inches(1.0)
            },
            "bullet_style": "●",
            "slide_number": True,
            "shape_roundness": 0.2,
        },
        # 기타 설정
        "settings": {
            "footer_text": "© 남양주시",
            "include_logo": False,
            "logo_path": None,
            "background_style": "solid",
            "accent_line": True,  # 상단에 액센트 라인 추가
        }
    }
}

# 기본 데이터 처리 및 추출 함수
def extract_text_from_pdf(pdf_file):
    """PDF 파일에서 텍스트 추출"""
    text = ""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
        temp_file.write(pdf_file.read())
        temp_path = temp_file.name
    
    try:
        with open(temp_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n\n"
    finally:
        os.unlink(temp_path)
    
    return text

def extract_text_from_docx(docx_file):
    """Word 파일에서 텍스트 추출"""
    text = ""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
        temp_file.write(docx_file.read())
        temp_path = temp_file.name
    
    try:
        doc = docx.Document(temp_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    finally:
        os.unlink(temp_path)
    
    return text

# AI 기반 슬라이드 분석 및 생성 관련 함수들

def enhance_with_openai(text, num_slides, api_key, temperature=0.7):
    """OpenAI를 사용하여 문서 구조 개선 및 슬라이드 구성"""
    try:
        client = OpenAI(api_key=api_key)
        
        # 시스템 프롬프트 정의
        system_prompt = """
        너는 문서를 분석하여 PowerPoint 슬라이드로 변환하는 전문가여야해!
        주어진 텍스트를 분석하고, 필요한 수의 슬라이드로 구성된 프레젠테이션 구조를 만들어줘
        각 슬라이드는 명확한 '제목'과 '내용'을 포함해야하고
        내용은 불릿 포인트 형태의 간결한 문장으로 구성해야해
        
        첫 번째 슬라이드는 전체 프레젠테이션의 제목 슬라이드로 구성하고.
        나머지 슬라이드는 내용을 논리적으로 구분하여 배치해
        절대로 주요항목이 빠지면 안 되고 필요하다면 슬라이드 수를 늘려!!
        아래 형식의 JSON으로 응답해줘:
        [
          {
            "title": "슬라이드 제목",
            "content": ["불릿 포인트 1", "불릿 포인트 2", "불릿 포인트 3"...]
          }
        ]
        """
        
        # 사용자 프롬프트 정의
        user_prompt = f"""
        다음 텍스트를 분석하고 반드시 {num_slides}개의 슬라이드로 구성된 프레젠테이션으로 변환해:
        
        {text}
        
        각 슬라이드는 잘 구성된 제목과 7-8개 이내의 핵심 불릿 포인트를 포함해야해.
        첫 번째 슬라이드는 표지/제목 슬라이드로 구성하고
        슬라이드가 8개 이상이되면 마지막 슬라이드는 요약이나 문의처 정보로 구성해줘
        반드시 JSON 형식으로 응답해야해!
        """
        
        # 간소화된 응답 요청 형식
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=temperature
        )
        
        # 응답 처리
        content = response.choices[0].message.content
        
        # 여러 JSON 파싱 시도
        try:
            # 1. 직접 JSON 파싱 시도
            slides_data = json.loads(content)
            
            # 2. 딕셔너리에서 slides 키 확인
            if isinstance(slides_data, dict) and "slides" in slides_data:
                slides_data = slides_data["slides"]
                
        except json.JSONDecodeError:
            # 3. 텍스트에서 JSON 부분 추출 시도
            json_pattern = r'(\[\s*\{.*\}\s*\])'
            json_match = re.search(json_pattern, content, re.DOTALL)
            
            if json_match:
                try:
                    slides_data = json.loads(json_match.group(1))
                except:
                    # 4. 마크다운 코드 블록에서 JSON 추출 시도
                    code_block_pattern = r'```(?:json)?\s*([\s\S]*?)```'
                    code_match = re.search(code_block_pattern, content, re.DOTALL)
                    if code_match:
                        try:
                            slides_data = json.loads(code_match.group(1))
                        except:
                            # 5. 실패하면 텍스트 파싱 시도
                            slides_data = extract_slides_from_text(content)
                    else:
                        slides_data = extract_slides_from_text(content)
            else:
                # 6. 모든 방법 실패 시 텍스트 파싱 시도
                slides_data = extract_slides_from_text(content)
        
        # 응답 슬라이드 데이터 확인 및 처리
        if not isinstance(slides_data, list):
            st.warning("AI가 올바른 형식의 슬라이드 데이터를 반환하지 않았습니다. 기본 분석을 사용합니다.")
            return fallback_parse_document(text, num_slides)
        
        # 각 슬라이드 내용 확인 및 정리
        processed_slides = []
        for slide in slides_data:
            if not isinstance(slide, dict):
                continue
                
            title = slide.get("title", "슬라이드")
            if not title:
                title = "슬라이드"
                
            # 내용 처리
            content = slide.get("content", [])
            if isinstance(content, str):
                # 줄바꿈으로 분할해 리스트로 변환
                content = [point.strip() for point in content.split('\n') if point.strip()]
                # 또는 문장 단위로 분할
                if len(content) <= 1 and content:
                    content = [s.strip() for s in re.split(r'[.!?]', content[0]) if s.strip()]
            
            processed_slides.append({
                "title": title,
                "content": content if content else ["내용이 필요합니다"]
            })
        
        if processed_slides:
            return processed_slides
        else:
            st.warning("슬라이드 구조를 생성하지 못했습니다. 기본 분석으로 진행합니다.")
            return fallback_parse_document(text, num_slides)
            
    except Exception as e:
        st.error(f"OpenAI 처리 중 오류 발생: {str(e)}")
        return fallback_parse_document(text, num_slides)

def enhance_with_gemini(text, num_slides, api_key, temperature=0.7):
    """Google Gemini를 사용하여 문서 구조 개선 및 슬라이드 구성"""
    try:
        genai.configure(api_key=api_key)
        
        prompt = f"""
        너는 문서를 분석하여 PowerPoint 슬라이드로 변환하는 전문가야
        다음 텍스트를 분석하고 {num_slides}개의 슬라이드로 구성된 프레젠테이션으로 변환해줘:
        
        {text}
        
        각 슬라이드는 잘 구성된 제목과 7-8개 이내의 핵심 불릿 포인트를 포함해야해
        첫 번째 슬라이드는 표지/제목 슬라이드로 구성하고
        슬라이드가 8개 이상이면 마지막 슬라이드는 요약이나 문의처 정보로 구성해주세요.
        
        아래 형식의 JSON으로 응답해주세요:
        [
          {{
            "title": "슬라이드 제목",
            "content": ["불릿 포인트 1", "불릿 포인트 2", "불릿 포인트 3"]
          }}
        ]
        
        반드시 정확한 JSON 형식만 응답하고, 추가 설명은 필요하지 않아아.
        """
        
        model = genai.GenerativeModel("gemini-2.0-flash-lite", 
                                     generation_config={"temperature": temperature})
        response = model.generate_content(prompt)
        
        # 응답 처리
        content = response.text
        
        # 여러 JSON 파싱 시도
        try:
            # 1. 직접 JSON 파싱 시도
            slides_data = json.loads(content)
            
            # 2. 딕셔너리에서 slides 키 확인
            if isinstance(slides_data, dict) and "slides" in slides_data:
                slides_data = slides_data["slides"]
                
        except json.JSONDecodeError:
            # 3. 텍스트에서 JSON 부분 추출 시도
            json_pattern = r'(\[\s*\{.*\}\s*\])'
            json_match = re.search(json_pattern, content, re.DOTALL)
            
            if json_match:
                try:
                    slides_data = json.loads(json_match.group(1))
                except:
                    # 4. 마크다운 코드 블록에서 JSON 추출 시도
                    code_block_pattern = r'```(?:json)?\s*([\s\S]*?)```'
                    code_match = re.search(code_block_pattern, content, re.DOTALL)
                    if code_match:
                        try:
                            slides_data = json.loads(code_match.group(1))
                        except:
                            # 5. 실패하면 텍스트 파싱 시도
                            slides_data = extract_slides_from_text(content)
                    else:
                        slides_data = extract_slides_from_text(content)
            else:
                # 6. 모든 방법 실패 시 텍스트 파싱 시도
                slides_data = extract_slides_from_text(content)
        
        # 응답 슬라이드 데이터 확인 및 처리
        if not isinstance(slides_data, list):
            st.warning("AI가 올바른 형식의 슬라이드 데이터를 반환하지 않았습니다. 기본 분석을 사용합니다.")
            return fallback_parse_document(text, num_slides)
        
        # 각 슬라이드 내용 확인 및 정리
        processed_slides = []
        for slide in slides_data:
            if not isinstance(slide, dict):
                continue
                
            title = slide.get("title", "슬라이드")
            if not title:
                title = "슬라이드"
                
            # 내용 처리
            content = slide.get("content", [])
            if isinstance(content, str):
                # 줄바꿈으로 분할해 리스트로 변환
                content = [point.strip() for point in content.split('\n') if point.strip()]
                # 또는 문장 단위로 분할
                if len(content) <= 1 and content:
                    content = [s.strip() for s in re.split(r'[.!?]', content[0]) if s.strip()]
            
            processed_slides.append({
                "title": title,
                "content": content if content else ["내용이 필요합니다"]
            })
        
        if processed_slides:
            return processed_slides
        else:
            st.warning("슬라이드 구조를 생성하지 못했습니다. 기본 분석으로 진행합니다.")
            return fallback_parse_document(text, num_slides)
            
    except Exception as e:
        st.error(f"Gemini 처리 중 오류 발생: {str(e)}")
        return fallback_parse_document(text, num_slides)

def extract_slides_from_text(text):
    """텍스트 응답에서 슬라이드 구조 추출"""
    slides = []
    
    # 슬라이드 구분 패턴 탐지
    slide_patterns = [
        r'(?:슬라이드|Slide)\s*\d+[:.]\s*(.*?)(?=(?:슬라이드|Slide)\s*\d+[:.]\s*|$)',  # 슬라이드 번호 형식
        r'#\s*(.*?)(?=\n#\s*|$)',  # 마크다운 헤딩 형식
        r'\*\*(.*?)\*\*\s*\n',  # 볼드 텍스트 제목 형식
    ]
    
    for pattern in slide_patterns:
        matches = re.finditer(pattern, text, re.DOTALL | re.MULTILINE)
        slide_titles = [match.group(1).strip() for match in matches if match.group(1).strip()]
        
        if slide_titles:
            # 제목이 발견되면 내용 추출 시도
            sections = re.split(pattern, text)[1:]  # 첫 부분 제외
            
            for i, title in enumerate(slide_titles):
                content = []
                
                if i < len(sections):
                    section_text = sections[i]
                    # 불릿 포인트 추출
                    bullet_patterns = [
                        r'[•\-\*]\s*(.*?)(?=\n[•\-\*]|$)',  # 기본 불릿 형식
                        r'\d+\.\s*(.*?)(?=\n\d+\.|$)',  # 숫자 리스트 형식
                    ]
                    
                    for bp in bullet_patterns:
                        bullet_matches = re.finditer(bp, section_text, re.MULTILINE)
                        points = [m.group(1).strip() for m in bullet_matches if m.group(1).strip()]
                        if points:
                            content.extend(points)
                    
                    # 불릿이 발견되지 않으면 문장 단위로 분할
                    if not content:
                        sentences = re.split(r'(?<=[.!?])\s+', section_text)
                        content = [s.strip() for s in sentences if s.strip()]
                
                if not content:
                    content = ["내용 없음"]
                
                slides.append({
                    "title": title,
                    "content": content
                })
            
            # 적어도 하나의 슬라이드가 생성되었으면 중단
            if slides:
                break
    
    # 패턴 매칭이 실패한 경우 기본 슬라이드 생성
    if not slides:
        slides = [
            {"title": "프레젠테이션", "content": ["자동 생성된 슬라이드"]},
        ]
    
    return slides

def fallback_parse_document(text, min_slides=5):
    """기본 문서 구조 파싱 (AI API 실패 시 대체 방법)"""
    # 문서 제목 추출 시도
    lines = text.split('\n')
    title = "프레젠테이션"
    for line in lines[:5]:  # 처음 5줄에서 제목 찾기
        if line.strip() and len(line.strip()) < 50:
            title = line.strip()
            break
    
    # 첫 번째 슬라이드 (제목 슬라이드)
    slides = [{
        "title": title,
        "content": ["자동 생성된 프레젠테이션"]
    }]
    
    # 섹션 구분 패턴 (■, -, □, ※ 등의 기호로 시작하는 줄)
    section_patterns = [
        r'^■\s*(.*?)$',  # ■ 기호로 시작
        r'^-\s*(.*?)$',   # - 기호로 시작
        r'^□\s*(.*?)$',  # □ 기호로 시작
        r'^※\s*(.*?)$',  # ※ 기호로 시작
        r'^\d+\.\s*(.*?)$'  # 숫자 + 점으로 시작
    ]
    
    # 텍스트에서 섹션 추출
    sections = []
    current_section = {"title": "개요", "content": []}
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        is_section_start = False
        for pattern in section_patterns:
            match = re.match(pattern, line)
            if match:
                # 이전 섹션 저장
                if current_section["content"]:
                    sections.append(current_section)
                
                # 새 섹션 시작
                section_title = match.group(1).strip()
                if not section_title:
                    section_title = "세부 내용"
                current_section = {"title": section_title, "content": []}
                is_section_start = True
                break
        
        if not is_section_start:
            # 현재 섹션에 내용 추가
            if line.startswith('-') or line.startswith('•'):
                line = line[1:].strip()
            current_section["content"].append(line)
    
    # 마지막 섹션 추가
    if current_section["content"]:
        sections.append(current_section)
    
    # 섹션이 너무 적으면 긴 섹션 분할
    while len(sections) < min_slides - 2:  # 제목 슬라이드와 마무리 슬라이드 제외
        longest_section = max(sections, key=lambda s: len(s["content"]))
        if len(longest_section["content"]) < 2:
            break
            
        # 섹션 분할
        mid = len(longest_section["content"]) // 2
        new_section = {
            "title": longest_section["title"] + " (계속)",
            "content": longest_section["content"][mid:]
        }
        longest_section["content"] = longest_section["content"][:mid]
        
        # 새 섹션 추가
        sections.append(new_section)
    
    # 섹션을 슬라이드로 변환
    for section in sections:
        slides.append({
            "title": section["title"],
            "content": section["content"]
        })
    
    # 마무리 슬라이드 추가
    contact_info = "추가 정보나 문의사항이 있으시면 연락해 주세요."
    for line in lines:
        if "문의처" in line or "연락처" in line or "☎" in line:
            contact_info = line.strip()
            break
    
    slides.append({
        "title": "감사합니다",
        "content": ["요약 및 문의처", contact_info]
    })
    
    # 슬라이드 수가 최소 요구사항을 충족하는지 확인
    while len(slides) < min_slides:
        # 더미 슬라이드 추가
        slides.append({
            "title": "추가 정보",
            "content": ["필요한 경우 추가 정보를 제공합니다."]
        })
    
    return slides

# PPT 생성 및 스타일링 함수들

def apply_background(slide, template):
    """슬라이드 배경 스타일 적용"""
    # 템플릿 정보에서 필요한 값 추출
    colors = template["colors"]
    bg_color = colors["bg_color"]
    background_style = template["settings"].get("background_style", "solid")
    
    # 배경색 설정 (기본)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*bg_color)
    
    # 그라데이션 배경
    if background_style == "gradient" and template["settings"].get("gradient_stops"):
        # 참고: python-pptx는 현재 그라데이션 배경을 직접 지원하지 않음
        # 그라데이션 배경을 적용하려면 CustomXML 조작이 필요함
        # 여기서는 단색으로 대체하고 실제 구현은 라이브러리의 향후 버전이나 확장 메서드가 필요
        pass
    
    # 패턴 배경
    elif background_style == "pattern" and template["settings"].get("pattern_style"):
        # 참고: 패턴 배경도 직접 지원하지 않음
        # 이것도 단색으로 대체
        pass
    
    # 이미지 배경
    elif background_style == "image" and template["settings"].get("background_image"):
        # 이미지 배경 적용 로직
        pass

def apply_text_style(text_frame, color, font_name, font_size, bold=False, italic=False, alignment=PP_ALIGN.LEFT):
    """텍스트 프레임에 스타일 적용"""
    if isinstance(text_frame, object) and hasattr(text_frame, 'paragraphs'):
        # 텍스트 프레임에 적용
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = alignment
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(*color)
                run.font.name = font_name
                run.font.size = font_size
                run.font.bold = bold
                run.font.italic = italic
    else:
        # 단일 단락에 적용
        text_frame.alignment = alignment
        for run in text_frame.runs:
            run.font.color.rgb = RGBColor(*color)
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = bold
            run.font.italic = italic

def customize_bullet(paragraph, bullet_style, bullet_color):
    """불릿 스타일 커스터마이징"""
    # 불릿 스타일은 현재 python-pptx에서 직접 지원하지 않음
    # 기본 불릿만 설정하고 스타일 변경은 지원하지 않음
    
    # 불릿 활성화
    if hasattr(paragraph, '_p') and hasattr(paragraph._p, 'get_or_add_pPr'):
        paragraph._p.get_or_add_pPr().set('bullet', '1')

def add_footer(slide, footer_text, footer_color, footer_font):
    """슬라이드에 푸터 추가"""
    # 푸터 텍스트 상자 추가
    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(9.0)
    height = Inches(0.3)
    
    footer_shape = slide.shapes.add_textbox(left, top, width, height)
    footer_text_frame = footer_shape.text_frame
    footer_text_frame.clear()
    
    p = footer_text_frame.add_paragraph()
    p.text = footer_text
    p.alignment = PP_ALIGN.CENTER
    
    for run in p.runs:
        run.font.name = footer_font["name"]
        run.font.size = footer_font["size"]
        run.font.color.rgb = RGBColor(*footer_color)
        run.font.bold = footer_font.get("bold", False)

def add_slide_number(slide, current_number, total_slides, footer_color, footer_font):
    """슬라이드 번호 추가"""
    # 슬라이드 번호 텍스트 상자 추가
    left = Inches(9.3)
    top = Inches(7.0)
    width = Inches(0.5)
    height = Inches(0.3)
    
    number_shape = slide.shapes.add_textbox(left, top, width, height)
    number_text_frame = number_shape.text_frame
    number_text_frame.clear()
    
    p = number_text_frame.add_paragraph()
    p.text = f"{current_number}/{total_slides}"
    p.alignment = PP_ALIGN.RIGHT
    
    for run in p.runs:
        run.font.name = footer_font["name"]
        run.font.size = footer_font["size"]
        run.font.color.rgb = RGBColor(*footer_color)

def add_logo(slide, logo_path):
    """슬라이드에 로고 추가"""
    # 이미지 파일이 존재하는지 확인
    if not os.path.exists(logo_path):
        return
    
    # 로고 이미지 추가
    left = Inches(9.0)
    top = Inches(0.3)
    width = Inches(0.8)
    height = Inches(0.4)
    
    # 이미지 추가 (실제 크기 비율 유지)
    slide.shapes.add_picture(logo_path, left, top, width, height)

def add_accent_line(slide, color, width=Inches(2), height=Inches(0.05)):
    """상단 액센트 라인 추가 (미니멀 디자인 등에서 사용)"""
    left = Inches(0.5)
    top = Inches(1.0)
    
    # 라인 추가
    shape_type = MSO_SHAPE.RECTANGLE
    line_shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(*color)
    line_shape.line.fill.background()  # 테두리 없음

def create_enhanced_ppt(slides, template_name="기본 템플릿"):
    """향상된 PPT 템플릿으로 생성"""
    prs = Presentation()
    
    # 선택한 템플릿 정보
    template = ENHANCED_TEMPLATES.get(template_name, ENHANCED_TEMPLATES["기본 템플릿"])
    
    # 색상 정보 추출
    colors = template["colors"]
    bg_color = colors["bg_color"]
    title_color = colors["title_color"]
    text_color = colors["text_color"]
    accent_color = colors["accent_color"]
    bullet_color = colors["bullet_color"]
    footer_color = colors["footer_color"]
    
    # 폰트 정보 추출
    fonts = template["fonts"]
    title_font = fonts["title_font"]
    subtitle_font = fonts["subtitle_font"]
    body_font = fonts["body_font"]
    footer_font = fonts["footer_font"]
    
    # 스타일 정보 추출
    styles = template["styles"]
    title_align = styles["title_align"]
    body_align = styles["body_align"]
    margins = styles.get("margins", {})
    bullet_style = styles.get("bullet_style", "•")
    slide_number = styles.get("slide_number", True)
    
    # 기타 설정 추출
    settings = template["settings"]
    footer_text = settings.get("footer_text", "© 2025 남양주시")
    include_logo = settings.get("include_logo", False)
    background_style = settings.get("background_style", "solid")
    
    # 슬라이드 레이아웃 가져오기
    title_slide_layout = prs.slide_layouts[0]  # 제목 슬라이드
    content_slide_layout = prs.slide_layouts[1]  # 제목 및 내용 슬라이드
    
    # 슬라이드 생성
    for i, slide_data in enumerate(slides):
        if i == 0:  # 첫 번째 슬라이드는 제목 슬라이드로 생성
            slide = prs.slides.add_slide(title_slide_layout)
            
            # 배경 스타일 적용
            apply_background(slide, template)
            
            # 제목 설정
            title = slide.shapes.title
            title.text = slide_data.get("title", "프레젠테이션")
            
            # 제목 텍스트 스타일 적용
            apply_text_style(title.text_frame, 
                            color=title_color,
                            font_name=title_font["name"],
                            font_size=title_font["size"],
                            bold=title_font["bold"],
                            italic=title_font["italic"],
                            alignment=title_align)
            
            # 부제목 설정
            if "content" in slide_data and slide_data["content"]:
                subtitle = slide.placeholders[1]
                if isinstance(slide_data["content"], list):
                    subtitle.text = slide_data["content"][0]
                else:
                    subtitle.text = str(slide_data["content"])
                
                # 부제목 텍스트 스타일 적용
                apply_text_style(subtitle.text_frame, 
                                color=accent_color,
                                font_name=subtitle_font["name"],
                                font_size=subtitle_font["size"],
                                bold=subtitle_font["bold"],
                                italic=subtitle_font["italic"],
                                alignment=title_align)
            
            # 제목 슬라이드에 액센트 라인 추가 (미니멀리즘 템플릿 등에서 사용)
            if settings.get("accent_line", False):
                add_accent_line(slide, accent_color, width=Inches(2), height=Inches(0.05))
            
        else:
            # 내용 슬라이드 생성
            slide = prs.slides.add_slide(content_slide_layout)
            
            # 배경 스타일 적용
            apply_background(slide, template)
            
            # 제목 설정
            title = slide.shapes.title
            title.text = slide_data.get("title", "슬라이드")
            
            # 제목 텍스트 스타일 적용
            apply_text_style(title.text_frame, 
                            color=title_color,
                            font_name=title_font["name"],
                            font_size=title_font["size"],
                            bold=title_font["bold"],
                            italic=title_font["italic"],
                            alignment=title_align)
            
            # 내용 설정
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.clear()
            
            # 콘텐츠가 리스트인 경우 (불릿 포인트)
            if "content" in slide_data:
                content = slide_data["content"]
                if isinstance(content, list):
                    for idx, point in enumerate(content):
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
                        
                        # 본문 텍스트 스타일 적용
                        apply_text_style(p, 
                                        color=text_color,
                                        font_name=body_font["name"],
                                        font_size=body_font["size"],
                                        bold=body_font["bold"],
                                        italic=body_font["italic"],
                                        alignment=body_align)
                        
                        # 불릿 스타일 커스터마이징
                        customize_bullet(p, bullet_style, bullet_color)
                else:
                    # 문자열인 경우
                    p = tf.add_paragraph()
                    p.text = str(content)
                    
                    # 본문 텍스트 스타일 적용
                    apply_text_style(p, 
                                    color=text_color,
                                    font_name=body_font["name"],
                                    font_size=body_font["size"],
                                    bold=body_font["bold"],
                                    italic=body_font["italic"],
                                    alignment=body_align)
        
        # 공통 요소 추가 (푸터, 로고, 슬라이드 번호 등)
        if i > 0 or settings.get("footer_on_title", False):  # 첫 슬라이드에는 보통 푸터를 넣지 않음 (설정 가능)
            # 푸터 추가
            add_footer(slide, footer_text, footer_color, footer_font)
            
            # 슬라이드 번호 추가
            if slide_number:
                add_slide_number(slide, i+1, len(slides), footer_color, footer_font)
            
            # 로고 추가 (설정된 경우)
            if include_logo and settings.get("logo_path"):
                add_logo(slide, settings["logo_path"])
    
    # 메모리 스트림에 PPT 저장
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output

# 메인 함수 - Streamlit UI 및 실행 로직
def run():
    st.title("📊 AI 기반 문서 PPT 변환기")
    st.caption("계획서나 보고서를 업로드하거나 텍스트를 입력하면 AI가 분석하여 PowerPoint로 변환해드립니다.")

    # 메인 영역 스타일 적용
    st.markdown("""
    <style>
    .content-output {
        background-color: #f0f0f0;
        border: 1px solid #ddd;
        border-radius: 8px;
        padding: 15px;
        margin-top: 10px;
        max-height: 600px;
        overflow-y: auto;
    }
    .output-title {
        font-size: 20px;
        font-weight: bold;
        margin-bottom: 10px;
        color: #2c3e50;
    }
    .output-text {
        font-size: 14px;
        line-height: 1.5;
    }
    h1 {
        font-size: 30px !important;
    }
    h3 {
        font-size: 16px !important;
        margin-top: 8px !important;
        margin-bottom: 4px !important;
    }
    /* 텍스트 영역 상단 여백 조정 */
    .stTextArea {
        margin-top: 0rem !important;
        margin-bottom: 2rem !important;        
    }
    /* 선택 옵션 스타일 */
    .stRadio label, .stSelectbox label {
        font-size: 14px !important;
    }
    /* 선택 컨테이너 스타일 */
    .option-container {
        background-color: #f8f9fa;
        padding: 10px;
        border-radius: 8px;
        margin-bottom: 15px;
    }
    /* 템플릿 미리보기 스타일 */
    .template-preview {
        border: 1px solid #ddd;
        border-radius: 5px;
        padding: 5px;
        margin-bottom: 10px;
    }
    .color-box {
        height: 20px;
        margin-bottom: 5px;
        border-radius: 3px;
        border: 1px solid #ddd;
    }
    </style>
    """, unsafe_allow_html=True)

    # API 키 확인
    if not OPENAI_API_KEY or not GEMINI_API_KEY:
        st.error("API 키가 설정되지 않았습니다. .env 파일에 OPENAI_API_KEY와 GEMINI_API_KEY를 설정해주세요.")
        st.stop()

    # 세션 상태 초기화
    if "document_text" not in st.session_state:
        st.session_state.document_text = ""
    if "ppt_generated" not in st.session_state:
        st.session_state.ppt_generated = False
    if "slides_preview" not in st.session_state:
        st.session_state.slides_preview = []

    # 사이드바 설정
    with st.sidebar:
        st.divider()
        # AI 모델 선택
        model_provider = st.radio(
            "🤖 AI모델 선택",
            ["OpenAI GPT-4o", "Google Gemini-2.0"]
        )
        
        # 공통 설정
        temperature = st.slider("⚙️ 창의성 수준", min_value=0.0, max_value=1.0, value=0.7, step=0.1)
        
        # 템플릿 선택 (향상된 템플릿 목록 사용)
        template_name = st.selectbox(
            "PPT 템플릿 선택",
            list(ENHANCED_TEMPLATES.keys()),
            index=0
        )
        
        # 템플릿 미리보기
        if st.checkbox("템플릿 미리보기", value=False):
            selected_template = ENHANCED_TEMPLATES[template_name]
            
            # 색상 팔레트 미리보기
            st.markdown("##### 색상 팔레트")
            colors = selected_template["colors"]
            
            # 색상 그리드 (2열로 표시)
            cols1, cols2 = st.columns(2)
            with cols1:
                st.markdown(f"<div class='color-box' style='background-color:rgb{colors['bg_color']};'></div>배경색", unsafe_allow_html=True)
                st.markdown(f"<div class='color-box' style='background-color:rgb{colors['title_color']};'></div>제목색", unsafe_allow_html=True)
                st.markdown(f"<div class='color-box' style='background-color:rgb{colors['text_color']};'></div>텍스트색", unsafe_allow_html=True)
            
            with cols2:
                st.markdown(f"<div class='color-box' style='background-color:rgb{colors['accent_color']};'></div>강조색", unsafe_allow_html=True)
                st.markdown(f"<div class='color-box' style='background-color:rgb{colors['bullet_color']};'></div>불릿색", unsafe_allow_html=True)
                st.markdown(f"<div class='color-box' style='background-color:rgb{colors['footer_color']};'></div>푸터색", unsafe_allow_html=True)
            
            # 폰트 정보
            st.markdown("##### 폰트 설정")
            fonts = selected_template["fonts"]
            st.caption(f"제목: {fonts['title_font']['name']} ({fonts['title_font']['size'].__str__().replace('Pt', '')}pt)")
            st.caption(f"본문: {fonts['body_font']['name']} ({fonts['body_font']['size'].__str__().replace('Pt', '')}pt)")
            
            # 스타일 정보
            st.markdown("##### 스타일 설정")
            styles = selected_template["styles"]
            st.caption(f"불릿 스타일: {styles['bullet_style']}")
            st.caption(f"슬라이드 번호: {'표시' if styles['slide_number'] else '숨김'}")
            
            # 기타 설정
            settings = selected_template["settings"]
            st.caption(f"푸터 텍스트: {settings['footer_text']}")
            st.caption(f"배경 스타일: {settings['background_style']}")
        
        # 슬라이드 수 설정
        num_slides = st.slider(
            "슬라이드 수",
            min_value=3,
            max_value=20,
            value=5,
            help="생성할 슬라이드의 개수를 설정합니다."
        )
        
        st.divider()
        st.caption("© 2025 남양주시 AI 문서 PPT 변환기")
    
    # 문서 입력 영역
    st.subheader("문서 업로드 또는 텍스트 입력")
    
    # 입력 방식 선택 (가로로 배치)
    input_method = st.radio("입력 방식 선택", ["문서 업로드", "텍스트 직접 입력"], horizontal=True)
    
    document_text = ""
    
    if input_method == "문서 업로드":
        uploaded_file = st.file_uploader("PDF 또는 Word 문서를 업로드하세요", type=["pdf", "docx"])
        
        if uploaded_file is not None:
            with st.spinner("문서에서 텍스트를 추출하는 중..."):
                # 파일 유형에 따라 텍스트 추출
                if uploaded_file.name.endswith('.pdf'):
                    document_text = extract_text_from_pdf(uploaded_file)
                elif uploaded_file.name.endswith('.docx'):
                    document_text = extract_text_from_docx(uploaded_file)
                
                if document_text:
                    st.session_state.document_text = document_text
                    st.success("문서 추출 성공!")
                    with st.expander("추출된 텍스트 보기"):
                        st.text_area("텍스트 내용", value=document_text, height=300)
                else:
                    st.error("문서에서 텍스트를 추출할 수 없습니다.")
    else:
        st.session_state.document_text = st.text_area(
            "변환할 텍스트를 입력하세요",
            value=st.session_state.document_text,
            height=300,
            placeholder="여기에 계획서나 보고서 내용을 입력하세요..."
        )
        document_text = st.session_state.document_text
    
    # PPT 변환 버튼
    if document_text and st.button("PPT 생성하기", type="primary", use_container_width=True):
        with st.spinner("PPT 파일 생성 중..."):
            # 진행 상황 표시
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            # 1단계: 문서 분석
            status_text.text("문서 분석 중...")
            progress_bar.progress(25)
            
            # 2단계: AI로 슬라이드 구조 생성
            status_text.text("AI로 문서 구조 분석 중...")
            
            # 선택한 AI 모델에 따라 처리
            if model_provider == "OpenAI GPT-4o":
                slides = enhance_with_openai(document_text, num_slides, OPENAI_API_KEY, temperature)
            else:  # Google Gemini
                slides = enhance_with_gemini(document_text, num_slides, GEMINI_API_KEY, temperature)
            
            # 세션에 슬라이드 미리보기 저장
            st.session_state.slides_preview = slides
            
            progress_bar.progress(50)
            
            # 3단계: 향상된 PPT 템플릿으로 생성
            status_text.text("PowerPoint 파일 생성 중...")
            ppt_file = create_enhanced_ppt(slides, template_name)
            progress_bar.progress(75)
            
            # 4단계: 완료
            status_text.text("완료!")
            progress_bar.progress(100)
            
            # 파일명 설정
            current_time = time.strftime("%Y%m%d_%H%M%S")
            filename = f"presentation_{current_time}.pptx"
            
            # 다운로드 버튼 표시
            st.success(f"PPT 생성이 완료되었습니다! ({len(slides)}개 슬라이드)")
            
            # 슬라이드 구조 미리보기
            with st.expander("슬라이드 구조 미리보기", expanded=True):
                for i, slide in enumerate(slides):
                    st.markdown(f"**슬라이드 {i+1}: {slide['title']}**")
                    if isinstance(slide['content'], list):
                        for point in slide['content']:
                            st.markdown(f"- {point}")
                    else:
                        st.text(str(slide['content']))
            
            # PPT 다운로드 버튼
            st.download_button(
                label="📥 PPT 파일 다운로드",
                data=ppt_file,
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
            
            # 세션 상태 업데이트
            st.session_state.ppt_generated = True
    
    # 사용 방법 안내
    st.markdown("---")
    st.markdown("""
    ### 💡 사용 방법
    
    1. **문서 업로드** 또는 **텍스트 직접 입력**을 선택하세요.
    2. PDF나 Word 문서를 업로드하거나 텍스트를 입력하세요.
    3. **PPT 생성하기** 버튼을 클릭하세요.
    4. 생성된 PPT 파일을 다운로드하세요.
    
    ### ✨ 특징
    
    - AI 기반 문서 분석 및 슬라이드 구성
    - 핵심 내용 자동 추출 및 불릿 포인트 생성
    - 다양한 디자인 템플릿 지원
    - 세련된 폰트와 색상 조합
    - 자동 푸터 및 슬라이드 번호 추가
    - PDF 및 Word 문서 지원
    """)

if __name__ == "__main__":
    run()